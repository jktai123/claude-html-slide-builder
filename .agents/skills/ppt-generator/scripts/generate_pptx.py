#!/usr/bin/env python3
"""
PPT 生成脚本 (Canva 質感 & 科技黃黑配色版本)

功能：根据 JSON 格式的数据生成 PowerPoint (.pptx) 文件
依赖：python-pptx>=0.6.21

使用方法：
    python generate_pptx.py --input ./ppt_data.json --output ./presentation.pptx

输入格式：详见 references/ppt_structure_guide.md
"""

import json
import argparse
import sys
from pathlib import Path
from typing import Dict, List, Any

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(f"错误：缺少依赖库 {e}")
    print("请安装依赖：pip install python-pptx>=0.6.21")
    sys.exit(1)


class PPTGenerator:
    """PPT 生成器类 (Canva 質感設計版本)"""

    # 全局主題色彩與字體常數
    BG_COLOR = RGBColor(18, 18, 18)        # 深灰色背景 (#121212)
    ACCENT_COLOR = RGBColor(255, 215, 0)   # 金黃強調色 (#FFD700)
    TEXT_COLOR = RGBColor(255, 255, 255)   # 白色文字 (#FFFFFF)
    MUTED_COLOR = RGBColor(142, 142, 147)  # 輔助灰色文字 (#8E8E93)
    CARD_BG_COLOR = RGBColor(30, 30, 30)   # 卡片背景深灰色 (#1E1E1E)
    CARD_BORDER_COLOR = RGBColor(45, 45, 48)# 卡片邊框中灰色 (#2D2D30)

    def __init__(self, input_json: Dict[str, Any]):
        """
        初始化生成器

        Args:
            input_json: 解析后的 JSON 数据
        """
        self.data = input_json
        self.prs = None
        self._validate_input()

    def _validate_input(self):
        """验证输入数据的合法性"""
        if 'metadata' not in self.data:
            raise ValueError("缺少 metadata 字段")

        if 'slides' not in self.data or not isinstance(self.data['slides'], list):
            raise ValueError("缺少 slides 字段或格式错误")

        if len(self.data['slides']) == 0:
            raise ValueError("slides 不能为空")

        # 验证每个幻灯片
        valid_layouts = {
            'TitleSlide', 'TitleAndContent', 'TwoColumnText', 'SectionHeader',
            'ContentWithCaption', 'BulletList', 'BlankSlide',
            'FlowCards', 'ComparisonTable', 'GeometryGrid'
        }

        for idx, slide in enumerate(self.data['slides']):
            if 'layout' not in slide:
                raise ValueError(f"幻灯片 {idx + 1} 缺少 layout 字段")

            if slide['layout'] not in valid_layouts:
                raise ValueError(
                    f"幻灯片 {idx + 1} 的 layout 值 '{slide['layout']}' 无效。"
                    f"有效值：{', '.join(valid_layouts)}"
                )

            if 'title' not in slide or not slide['title']:
                raise ValueError(f"幻灯片 {idx + 1} 缺少 title 字段或为空")

            if 'content' not in slide or not isinstance(slide['content'], list):
                raise ValueError(f"幻灯片 {idx + 1} 的 content 字段必须是数组")

    def generate(self) -> Presentation:
        """
        生成 PPT 对象

        Returns:
            Presentation: 生成的 PPT 对象
        """
        # 创建演示文稿对象
        self.prs = Presentation()

        # 設定 16:9 寬螢幕尺寸 (Canva / Keynote 黃金比例)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # 设置元数据
        self._set_metadata()

        # 生成所有幻灯片
        for slide_data in self.data['slides']:
            self._add_slide(slide_data)

        # 自我查核與版面防溢出驗證
        self.verify_canvas_safety()

        return self.prs

    def _set_metadata(self):
        """设置演示文稿元数据"""
        metadata = self.data.get('metadata', {})

        # 设置核心属性
        if 'title' in metadata:
            self.prs.core_properties.title = metadata['title']

        if 'author' in metadata:
            self.prs.core_properties.author = metadata['author']

        if 'subject' in metadata:
            self.prs.core_properties.subject = metadata['subject']

        if 'keywords' in metadata:
            self.prs.core_properties.keywords = metadata['keywords']

    def _add_slide(self, slide_data: Dict[str, Any]):
        """
        添加一张幻灯片
        """
        layout_type = slide_data['layout']
        title = slide_data['title']
        content = slide_data.get('content', [])
        notes = slide_data.get('notes', '')

        # 使用空白版面 (Layout index 6) 提供完全客製化幾何排版
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 設定純色暗黑背景
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.BG_COLOR

        # 繪製標題與幾何裝飾 (BlankSlide 與 TitleSlide 除外)
        if layout_type not in ['TitleSlide', 'BlankSlide']:
            self._draw_geometry_decor(slide)
            self._add_textbox_with_autofit(
                slide, Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.8),
                title, 36, self.ACCENT_COLOR, bold=True
            )

        # 依布局繪製卡片
        if layout_type == 'TitleSlide':
            self._render_title_slide(slide, title, content)
        elif layout_type == 'SectionHeader':
            self._render_section_header(slide, title, content)
        elif layout_type == 'TwoColumnText':
            self._render_two_column(slide, content)
        elif layout_type == 'FlowCards':
            self._render_flow_cards(slide, content)
        elif layout_type == 'ComparisonTable':
            self._render_comparison_table(slide, content)
        elif layout_type == 'GeometryGrid':
            self._render_geometry_grid(slide, content)
        elif layout_type == 'BulletList':
            self._render_bullet_list(slide, content)
        elif layout_type == 'ContentWithCaption':
            self._render_content_with_caption(slide, content)
        elif layout_type == 'TitleAndContent':
            self._render_title_and_content(slide, content)

        # 设置备注
        if notes:
            self._set_notes(slide, notes)

    def _draw_geometry_decor(self, slide):
        """繪製科技感幾何修飾線"""
        # 左上角黃色幾何修飾條
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(0.4), Inches(1.5), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.ACCENT_COLOR
        bar.line.fill.background()

    def _add_textbox_with_autofit(self, slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
        """字數溢出自動縮小防重疊文字框"""
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        # 防溢出與自動縮放
        char_count = len(str(text))
        adjusted_pt = font_size
        if char_count > 150:
            adjusted_pt = max(12, int(font_size * 0.65))
        elif char_count > 80:
            adjusted_pt = max(14, int(font_size * 0.8))

        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.name = 'Arial'
        p.font.size = Pt(adjusted_pt)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return tx_box

    def _add_multiline_textbox(self, slide, left, top, width, height, lines, font_size, color, bold=False, space_after=8):
        """多行自適應文字框"""
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        total_chars = sum(len(str(line)) for line in lines)
        adjusted_pt = font_size
        if total_chars > 250:
            adjusted_pt = max(12, int(font_size * 0.7))
        elif total_chars > 120:
            adjusted_pt = max(14, int(font_size * 0.82))

        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = str(line)
            p.font.name = 'Arial'
            p.font.size = Pt(adjusted_pt)
            p.font.color.rgb = color
            p.font.bold = bold
            p.space_after = Pt(space_after)
        return tx_box

    # --- 佈局渲染器 ---

    def _render_title_slide(self, slide, title, content):
        """封面頁：Canva 質感大版面標題"""
        # 左下側大標題
        self._add_textbox_with_autofit(
            slide, Inches(1.0), Inches(2.4), Inches(11.333), Inches(1.5),
            title, 52, self.ACCENT_COLOR, bold=True
        )

        # 副標題
        if content:
            self._add_multiline_textbox(
                slide, Inches(1.0), Inches(4.0), Inches(11.333), Inches(2.0),
                content, 22, self.TEXT_COLOR, space_after=10
            )

        # 右下角抽象幾何三角色塊
        triangle = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10.5), Inches(4.7), Inches(2.833), Inches(2.8))
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = self.ACCENT_COLOR
        triangle.line.fill.background()
        triangle.rotation = 180

    def _render_section_header(self, slide, title, content):
        """章節頁：高對比卡片"""
        # 居中大圓角卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(2.2), Inches(9.333), Inches(3.0))
        card.fill.solid()
        card.fill.fore_color.rgb = self.CARD_BG_COLOR
        card.line.color.rgb = self.ACCENT_COLOR
        card.line.width = Pt(2.0)

        # 章節標題
        self._add_textbox_with_autofit(
            slide, Inches(2.2), Inches(2.6), Inches(8.933), Inches(1.0),
            title, 38, self.ACCENT_COLOR, bold=True, align=PP_ALIGN.CENTER
        )

        # 描述
        if content:
            self._add_textbox_with_autofit(
                slide, Inches(2.2), Inches(3.8), Inches(8.933), Inches(1.0),
                "\n".join(content), 18, self.TEXT_COLOR, align=PP_ALIGN.CENTER
            )

    def _render_title_and_content(self, slide, content):
        """標準內容頁：左側文字卡片 + 右側幾何視覺框"""
        # 左文字卡片
        left_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(7.2), Inches(4.8))
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = self.CARD_BG_COLOR
        left_box.line.color.rgb = self.CARD_BORDER_COLOR

        text_lines = [str(item) for item in content if not (str(item).startswith('[') and str(item).endswith(']'))]
        image_annotations = [str(item) for item in content if str(item).startswith('[') and str(item).endswith(']')]

        self._add_multiline_textbox(
            slide, Inches(1.3), Inches(2.1), Inches(6.6), Inches(4.2),
            text_lines, 20, self.TEXT_COLOR, space_after=12
        )

        # 右幾何展示框
        right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.5), Inches(4.8))
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = self.CARD_BG_COLOR
        right_box.line.color.rgb = self.ACCENT_COLOR
        right_box.line.width = Pt(1.5)

        label = "幾何科技視覺設計"
        if image_annotations:
            label = image_annotations[0].strip('[]')

        self._add_textbox_with_autofit(
            slide, Inches(9.0), Inches(3.6), Inches(3.1), Inches(1.5),
            label, 18, self.ACCENT_COLOR, bold=True, align=PP_ALIGN.CENTER
        )

        # 頂部裝飾菱形
        decor = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(10.25), Inches(2.6), Inches(0.6), Inches(0.6))
        decor.fill.solid()
        decor.fill.fore_color.rgb = self.ACCENT_COLOR
        decor.line.fill.background()

    def _render_bullet_list(self, slide, content):
        """項目列表：橫向跑馬燈卡片排版"""
        N = len(content)
        if N == 0:
            return

        card_height = min(1.2, (5.0 - (N - 1) * 0.25) / N)

        for idx, item in enumerate(content):
            top = 1.8 + idx * (card_height + 0.25)

            # 橫向卡片
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(top), Inches(11.333), Inches(card_height))
            card.fill.solid()
            card.fill.fore_color.rgb = self.CARD_BG_COLOR
            card.line.color.rgb = self.CARD_BORDER_COLOR

            # 黃色序號標記塊
            b_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(top + (card_height - 0.4) / 2), Inches(0.4), Inches(0.4))
            b_box.fill.solid()
            b_box.fill.fore_color.rgb = self.ACCENT_COLOR
            b_box.line.fill.background()

            tf = b_box.text_frame
            tf.text = str(idx + 1)
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

            # 卡片內文
            self._add_textbox_with_autofit(
                slide, Inches(2.0), Inches(top + (card_height - 0.6) / 2), Inches(10.0), Inches(0.6),
                str(item), 18, self.TEXT_COLOR
            )

    def _render_two_column(self, slide, content):
        """雙欄卡片排版"""
        left_content = []
        right_content = []
        current_side = 'left'

        for item in content:
            item_str = str(item).strip()
            if '【右】' in item_str or item_str.startswith('[右'):
                current_side = 'right'
                cleaned = item_str.replace('【右】', '').replace('[右栏]', '').strip()
                if cleaned: right_content.append(cleaned)
            elif '【左】' in item_str or item_str.startswith('[左'):
                current_side = 'left'
                cleaned = item_str.replace('【左】', '').replace('[左栏]', '').strip()
                if cleaned: left_content.append(cleaned)
            else:
                if current_side == 'left':
                    left_content.append(item_str)
                else:
                    right_content.append(item_str)

        # 左欄卡片
        left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(5.4), Inches(4.8))
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = self.CARD_BG_COLOR
        left_box.line.color.rgb = self.ACCENT_COLOR
        left_box.line.width = Pt(1.5)

        # 右欄卡片
        right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.8), Inches(5.4), Inches(4.8))
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = self.CARD_BG_COLOR
        right_box.line.color.rgb = self.CARD_BORDER_COLOR

        self._add_multiline_textbox(slide, Inches(1.3), Inches(2.1), Inches(4.8), Inches(4.2), left_content, 18, self.TEXT_COLOR)
        self._add_multiline_textbox(slide, Inches(7.233), Inches(2.1), Inches(4.8), Inches(4.2), right_content, 18, self.TEXT_COLOR)

    def _render_flow_cards(self, slide, content):
        """流程卡布局 (FlowCards)"""
        N = len(content)
        if N == 0:
            return

        total_w = 11.333
        spacing = 0.4
        card_w = (total_w - (N - 1) * spacing) / N

        for idx, item in enumerate(content):
            left = 1.0 + idx * (card_w + spacing)

            # 步驟卡
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.2), Inches(card_w), Inches(4.4))
            card.fill.solid()
            card.fill.fore_color.rgb = self.CARD_BG_COLOR
            card.line.color.rgb = self.CARD_BORDER_COLOR

            # 數字標題
            self._add_textbox_with_autofit(
                slide, Inches(left + 0.2), Inches(2.4), Inches(card_w - 0.4), Inches(0.5),
                f"0{idx + 1}", 26, self.ACCENT_COLOR, bold=True, align=PP_ALIGN.CENTER
            )

            # 步驟說明
            self._add_textbox_with_autofit(
                slide, Inches(left + 0.2), Inches(3.1), Inches(card_w - 0.4), Inches(3.2),
                str(item), 16, self.TEXT_COLOR, align=PP_ALIGN.CENTER
            )

            # 連接箭頭
            if idx < N - 1:
                a_left = left + card_w + 0.05
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(a_left), Inches(4.2), Inches(0.3), Inches(0.25))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.ACCENT_COLOR
                arrow.line.fill.background()

    def _render_comparison_table(self, slide, content):
        """比較表排版 (ComparisonTable)"""
        table_data = []
        for line in content:
            parts = [p.strip() for p in line.split('|')]
            table_data.append(parts)

        if not table_data:
            return

        rows = len(table_data)
        cols = max(len(row) for row in table_data)

        # 建立 PPTX 表格
        table_shape = slide.shapes.add_table(rows, cols, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
        table = table_shape.table

        for r_idx, row in enumerate(table_data):
            for c_idx, val in enumerate(row):
                if c_idx >= len(row):
                    continue
                cell = table.cell(r_idx, c_idx)
                cell.fill.solid()

                if r_idx == 0:
                    # 表頭：黃底黑字
                    cell.fill.fore_color.rgb = self.ACCENT_COLOR
                    color = RGBColor(0, 0, 0)
                    bold = True
                else:
                    # 表身：深灰底白字
                    cell.fill.fore_color.rgb = self.CARD_BG_COLOR
                    color = self.TEXT_COLOR
                    bold = False

                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.name = 'Arial'
                    p.font.size = Pt(16)
                    p.font.color.rgb = color
                    p.font.bold = bold

    def _render_geometry_grid(self, slide, content):
        """幾何網格矩陣 (GeometryGrid)"""
        N = len(content)
        if N == 0:
            return

        points = []
        for item in content:
            if '::' in str(item):
                hdr, desc = str(item).split('::', 1)
                points.append((hdr.strip(), desc.strip()))
            else:
                points.append((str(item), ""))

        # 佈局分配
        if N <= 2:
            card_w, card_h = 5.4, 4.4
            coords = [(1.0, 2.2), (6.933, 2.2)]
        elif N == 3:
            card_w, card_h = 3.5, 4.4
            coords = [(1.0, 2.2), (4.9, 2.2), (8.8, 2.2)]
        else:
            card_w, card_h = 5.4, 2.05
            coords = [(1.0, 2.2), (6.933, 2.2), (1.0, 4.55), (6.933, 4.55)]

        for idx, (hdr, desc) in enumerate(points):
            if idx >= len(coords):
                break
            x, y = coords[idx]

            # 網格卡片
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = self.CARD_BG_COLOR
            card.line.color.rgb = self.CARD_BORDER_COLOR

            # 左上小黃點裝飾
            decor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.25), Inches(y + 0.25), Inches(0.12), Inches(0.12))
            decor.fill.solid()
            decor.fill.fore_color.rgb = self.ACCENT_COLOR
            decor.line.fill.background()

            self._add_textbox_with_autofit(
                slide, Inches(x + 0.5), Inches(y + 0.15), Inches(card_w - 0.8), Inches(0.5),
                hdr, 18, self.ACCENT_COLOR, bold=True
            )

            if desc:
                self._add_textbox_with_autofit(
                    slide, Inches(x + 0.5), Inches(y + 0.7), Inches(card_w - 0.8), Inches(card_h - 0.9),
                    desc, 14, self.TEXT_COLOR
                )

    def _render_content_with_caption(self, slide, content):
        """帶說明的視覺卡片"""
        left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(6.0), Inches(4.8))
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = self.CARD_BG_COLOR
        left_box.line.color.rgb = self.ACCENT_COLOR
        left_box.line.width = Pt(1.5)

        text_lines = []
        v_label = "科技視覺呈現區"
        for item in content:
            item_str = str(item).strip()
            if item_str.startswith('[') and item_str.endswith(']'):
                v_label = item_str.strip('[]')
            else:
                text_lines.append(item_str)

        # 幾何菱形裝飾
        decor = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(3.7), Inches(3.2), Inches(0.6), Inches(0.6))
        decor.fill.solid()
        decor.fill.fore_color.rgb = self.ACCENT_COLOR
        decor.line.fill.background()

        self._add_textbox_with_autofit(
            slide, Inches(1.2), Inches(4.0), Inches(5.6), Inches(1.0),
            v_label, 18, self.ACCENT_COLOR, bold=True, align=PP_ALIGN.CENTER
        )

        # 右側說明文字
        right_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(1.8), Inches(4.833), Inches(4.8))
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = self.CARD_BG_COLOR
        right_box.line.color.rgb = self.CARD_BORDER_COLOR

        self._add_multiline_textbox(slide, Inches(7.8), Inches(2.1), Inches(4.233), Inches(4.2), text_lines, 18, self.TEXT_COLOR)

    # --- 自我查核與安全驗證 ---

    def verify_canvas_safety(self):
        """自我查核：確認元件無重疊、無溢出畫布"""
        print("▶ 正在進行簡報自我查核與安全性驗證...")
        canvas_width = Inches(13.333)
        canvas_height = Inches(7.5)

        overlap_found = False
        overflow_found = False

        for s_idx, slide in enumerate(self.prs.slides):
            shapes = list(slide.shapes)
            # 1. 溢出邊界檢測
            for shape in shapes:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                if left < 0 or top < 0 or (left + width) > canvas_width or (top + height) > canvas_height:
                    print(f"  ⚠️ [溢出警告] 投影片 {s_idx+1} 的元件 '{shape.name}' 超出畫布邊界！"
                          f" (X: {left.inches:.2f}, Y: {top.inches:.2f}, W: {width.inches:.2f}, H: {height.inches:.2f})")
                    overflow_found = True

            # 2. 局部重疊交叉檢測
            for i in range(len(shapes)):
                for j in range(i + 1, len(shapes)):
                    s1 = shapes[i]
                    s2 = shapes[j]

                    if s1.width >= canvas_width or s2.width >= canvas_width:
                        continue
                    if s1.height < Inches(0.2) or s2.height < Inches(0.2):  # 忽略極細裝飾線
                        continue

                    left1, top1 = s1.left, s1.top
                    right1, bottom1 = left1 + s1.width, top1 + s1.height

                    left2, top2 = s2.left, s2.top
                    right2, bottom2 = left2 + s2.width, top2 + s2.height

                    x_overlap = max(0, min(right1, right2) - max(left1, left2))
                    y_overlap = max(0, min(bottom1, bottom2) - max(top1, top2))

                    if x_overlap > Inches(0.15) and y_overlap > Inches(0.15):
                        # 卡片與卡片內部文字的正常層疊不報警告
                        is_card_and_text = False
                        if (s1.shape_type == 1 and s2.has_text_frame) or (s2.shape_type == 1 and s1.has_text_frame):
                            is_card_and_text = True
                        if not is_card_and_text:
                            print(f"  ⚠️ [重疊警告] 投影片 {s_idx+1} 的元件 '{s1.name}' 與 '{s2.name}' 發生潛在重疊！")
                            overlap_found = True

        if not overlap_found and not overflow_found:
            print("  ✓ 自我查核完成：未發現重疊或溢出，版面安全符合 Canva 質感標準！")
        else:
            print("  ⚠️ 自我查核提示：檢測到部分版面安全提示，請手動確認生成效果。")

    def save(self, output_path: str):
        """
        保存 PPT 文件

        Args:
            output_path: 输出文件路径
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(output_path))
        print(f"✓ PPT 文件已生成: {output_path.absolute()}")


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description='根据 JSON 数据生成 PowerPoint 文件')
    parser.add_argument('--input', '-i', required=True, help='输入 JSON 文件路径')
    parser.add_argument('--output', '-o', required=True, help='输出 PPT 文件路径')

    args = parser.parse_args()

    # 读取输入文件
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"错误：输入文件不存在: {input_path}")
        sys.exit(1)

    print(f"正在读取输入文件: {input_path}")

    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"错误：JSON 解析失败 - {e}")
        sys.exit(1)
    except Exception as e:
        print(f"错误：读取文件失败 - {e}")
        sys.exit(1)

    # 生成 PPT
    print("正在生成 PPT...")

    try:
        generator = PPTGenerator(data)
        generator.generate()
        generator.save(args.output)
        print("✓ 生成成功！")
    except ValueError as e:
        print(f"错误：数据验证失败 - {e}")
        sys.exit(1)
    except Exception as e:
        print(f"错误：生成失败 - {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
